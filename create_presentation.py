#!/usr/bin/env python3
"""Create PowerPoint presentation summarizing the TEP ML Benchmarking project."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create presentation with widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Define colors
BLUE = RGBColor(0, 112, 192)
DARK_BLUE = RGBColor(0, 51, 102)
GREEN = RGBColor(0, 176, 80)
ORANGE = RGBColor(255, 153, 0)
RED = RGBColor(255, 0, 0)


def add_title_slide(title, subtitle=None):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(title, bullet_points, sub_bullets=None):
    """Add a content slide with title and bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.7))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(22)
        p.space_after = Pt(12)
        p.level = 0

        # Add sub-bullets if provided
        if sub_bullets and i in sub_bullets:
            for sub in sub_bullets[i]:
                sp = tf.add_paragraph()
                sp.text = f"    - {sub}"
                sp.font.size = Pt(18)
                sp.font.color.rgb = RGBColor(80, 80, 80)
                sp.space_after = Pt(6)
                sp.level = 1

    return slide


def add_table_slide(title, headers, rows, col_widths=None):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Calculate table dimensions
    num_rows = len(rows) + 1
    num_cols = len(headers)

    if col_widths is None:
        col_widths = [12.333 / num_cols] * num_cols

    table_width = sum(col_widths)
    table_left = (13.333 - table_width) / 2

    # Create table
    table = slide.shapes.add_table(num_rows, num_cols, Inches(table_left), Inches(1.4),
                                   Inches(table_width), Inches(0.5 * num_rows)).table

    # Set column widths
    for i, width in enumerate(col_widths):
        table.columns[i].width = Inches(width)

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(13)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)

    return slide


def add_two_column_slide(title, left_content, right_content):
    """Add a slide with two columns."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.9), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(18)
        p.space_after = Pt(8)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.3), Inches(5.9), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(18)
        p.space_after = Pt(8)

    return slide


# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
add_title_slide(
    "Benchmarking ML Anomaly Detection Methods\non the Tennessee Eastman Process",
    "Sunshine, Lyu, Botcha, Kulkarni, Pagaria, Alves, Kitchin\nDepartment of Chemical Engineering, Carnegie Mellon University"
)

# ============================================================================
# SLIDE 2: Motivation
# ============================================================================
add_content_slide(
    "Motivation",
    [
        "Industrial processes require reliable fault detection and diagnosis (FDD)",
        "Early fault detection prevents equipment damage, safety hazards, and costly downtime",
        "Tennessee Eastman Process (TEP) is a widely-used benchmark in process control",
        "Many ML methods proposed, but fair comparison is challenging:",
    ],
    {3: ["Different dataset splits", "Varying evaluation metrics", "No standardized benchmark"]}
)

# ============================================================================
# SLIDE 3: Research Objectives
# ============================================================================
add_content_slide(
    "Research Objectives",
    [
        "Create a rigorous, reproducible benchmark for ML-based fault detection",
        "Compare classical ML (XGBoost) vs. deep learning approaches",
        "Evaluate supervised multiclass classification and unsupervised anomaly detection",
        "Test temporal post-processing methods (HMM filtering)",
        "Assess generalization to new, unseen process data"
    ]
)

# ============================================================================
# SLIDE 4: Tennessee Eastman Process
# ============================================================================
add_content_slide(
    "Tennessee Eastman Process (TEP)",
    [
        "Realistic simulation of a chemical plant with 5 major unit operations",
        "52 process variables: 41 measurements + 11 manipulated variables",
        "17 fault types included (3 subtle faults excluded: 3, 9, 15)",
        "Fault types include: step changes, random variations, slow drifts, sticking valves",
        "Data from Harvard Dataverse (Rieth et al., 2017)"
    ]
)

# ============================================================================
# SLIDE 5: Dataset Construction
# ============================================================================
add_table_slide(
    "Dataset Construction",
    ["Dataset", "Samples", "Classes", "Purpose"],
    [
        ["Multiclass Train", "864,000", "18", "Model training"],
        ["Multiclass Val", "432,000", "18", "Hyperparameter tuning"],
        ["Multiclass Test", "2,880,000", "18", "Performance evaluation"],
        ["Binary Train", "50,000", "1 (normal)", "Autoencoder training"],
        ["Binary Test", "795,200", "2", "Anomaly detection eval"],
    ],
    col_widths=[3.0, 2.5, 2.5, 4.0]
)

# ============================================================================
# SLIDE 6: Dataset Key Features
# ============================================================================
add_content_slide(
    "Dataset Key Features",
    [
        "Perfectly balanced classes (standard deviation = 0.00)",
        "Zero data leakage: train/test from independent simulation runs",
        "Different fault introduction times (sample 21 vs 161)",
        "Strong statistical power: 95% CI width = +/-0.5%",
        "160,000 test samples per class (160x larger than MNIST per class)"
    ]
)

# ============================================================================
# SLIDE 7: Models Evaluated
# ============================================================================
add_two_column_slide(
    "Models Evaluated",
    [
        "Supervised Multiclass:",
        "  XGBoost (gradient boosting)",
        "  LSTM (recurrent network)",
        "  LSTM-FCN (LSTM + 1D CNN)",
        "  CNN-Transformer (attention)",
        "  TransKal (Transformer + Kalman)"
    ],
    [
        "Unsupervised Binary:",
        "  LSTM Autoencoder",
        "  Convolutional Autoencoder",
        "",
        "Post-processing:",
        "  HMM Classification Filter"
    ]
)

# ============================================================================
# SLIDE 8: 10-Series: Hyperparameter Tuning
# ============================================================================
add_content_slide(
    "10-Series: Hyperparameter Tuning",
    [
        "Optuna optimization with 50 trials per model",
        "Objective: Maximize weighted F1 score on validation set",
        "50% training data sampling (preserving run structure)",
        "MedianPruner for early stopping of unpromising trials",
        "Total tuning time: ~49 hours"
    ]
)

# ============================================================================
# SLIDE 9: Hyperparameter Tuning Results
# ============================================================================
add_table_slide(
    "Hyperparameter Tuning Results (10-Series)",
    ["Model", "Best Val F1", "Tuning Time", "Key Finding"],
    [
        ["XGBoost", "0.9236", "7h 28m", "max_depth=6, lr=0.20"],
        ["LSTM", "0.9878", "6h 37m", "seq_len=39, layers=3"],
        ["LSTM-FCN", "0.9896", "11h 1m", "seq_len=40, dropout=0.44"],
        ["CNN-Transformer", "0.9890", "8h 22m", "d_model=32, 4 heads"],
        ["TransKal", "0.9878", "12h 6m", "Kalman Q=5e-5, R=0.02"],
        ["LSTM-AE", "0.9817", "2h 34m", "threshold=98%"],
        ["Conv-AE", "0.9937", "1h 11m", "No transformer needed"],
    ],
    col_widths=[2.5, 2.0, 2.0, 5.5]
)

# ============================================================================
# SLIDE 10: Key Tuning Insights
# ============================================================================
add_content_slide(
    "Key Hyperparameter Insights",
    [
        "Optimal sequence length: 38-40 samples for all deep learning models",
        "Simpler architectures often outperformed complex ones",
        "Dropout varies significantly: LSTM (0.07) vs LSTM-FCN (0.44)",
        "Autoencoder threshold: ~98th percentile optimal for anomaly detection",
        "Conv-Autoencoder: Pure CNN beats CNN+Transformer hybrid"
    ]
)

# ============================================================================
# SLIDE 11: 20-Series: Final Training Results
# ============================================================================
add_table_slide(
    "Final Model Training Results (20-Series)",
    ["Model", "Test Accuracy", "F1 (weighted)", "Training Time"],
    [
        ["LSTM-FCN", "99.37%", "0.9937", "26m"],
        ["CNN-Transformer", "99.20%", "0.9920", "68m"],
        ["LSTM", "99.14%", "0.9914", "20m"],
        ["TransKal", "99.09%", "0.9909", "32m"],
        ["XGBoost", "93.91%", "0.9416", "40m"],
        ["Conv-AE (binary)", "99.39%", "0.9939", "<1m"],
        ["LSTM-AE (binary)", "96.96%", "0.9705", "12m"],
    ],
    col_widths=[3.0, 2.5, 2.5, 2.5]
)

# ============================================================================
# SLIDE 12: Training Observations
# ============================================================================
add_content_slide(
    "Training Observations (20-Series)",
    [
        "Deep learning models (99%+) significantly outperform XGBoost (94%)",
        "LSTM-FCN achieves best multiclass accuracy (99.37%)",
        "Temporal models excel: sequence-based > independent samples",
        "Some faults perfectly classified: Classes 6, 7, 14 achieve F1 = 1.0",
        "Challenging classes: Normal (0), Class 12, Class 18 have lower F1"
    ]
)

# ============================================================================
# SLIDE 13: 30-Series: HMM Filtering
# ============================================================================
add_content_slide(
    "30-Series: HMM Classification Filter",
    [
        "Post-processing to exploit temporal consistency in predictions",
        "Uses confusion matrix as emission model P(observation | true_class)",
        "\"Sticky\" transition model assumes faults persist over time",
        "Stickiness parameter gamma tested: 0.70, 0.85, 0.95",
        "Resets at beginning of each simulation run"
    ]
)

# ============================================================================
# SLIDE 14: HMM Filter Results
# ============================================================================
add_table_slide(
    "HMM Filter Results (30-Series)",
    ["Model", "Raw Acc", "Best gamma", "Filtered Acc", "Error Reduction"],
    [
        ["XGBoost", "93.91%", "0.95", "95.90%", "32.7%"],
        ["LSTM", "99.14%", "0.95", "99.14%", "0.7%"],
        ["LSTM-FCN", "99.37%", "0.70", "99.37%", "0.1%"],
        ["CNN-Transformer", "99.20%", "0.95", "99.21%", "1.2%"],
        ["TransKal", "99.08%", "0.85", "99.08%", "0.4%"],
    ],
    col_widths=[3.0, 2.0, 1.5, 2.5, 2.5]
)

# ============================================================================
# SLIDE 15: HMM Key Finding
# ============================================================================
add_content_slide(
    "HMM Filter: Key Finding",
    [
        "XGBoost benefits significantly: 32.7% error reduction (+2% accuracy)",
        "Deep learning models: <1.5% error reduction (minimal benefit)",
        "Reason: Deep learning already captures temporal patterns",
        "XGBoost treats samples independently -> needs post-hoc smoothing",
        "Recommendation: Use HMM for XGBoost; skip for neural networks"
    ]
)

# ============================================================================
# SLIDE 16: 40-Series: Generalization Testing
# ============================================================================
add_content_slide(
    "40-Series: New Data Evaluation",
    [
        "Critical test: Do models generalize to completely unseen data?",
        "Generated new dataset with tep-sim using different random seeds",
        "Same fault scenarios, same simulation parameters",
        "~2M new multiclass samples, ~1.2M new binary samples",
        "Tests real-world deployment scenario"
    ]
)

# ============================================================================
# SLIDE 17: Generalization Results
# ============================================================================
add_table_slide(
    "Generalization Results (40-Series)",
    ["Model", "Original Acc", "New Data Acc", "Change", "Assessment"],
    [
        ["CNN-Transformer", "99.20%", "99.57%", "+0.38%", "Excellent"],
        ["LSTM-FCN", "99.37%", "98.93%", "-0.44%", "Good"],
        ["LSTM", "99.14%", "98.91%", "-0.23%", "Good"],
        ["XGBoost", "93.91%", "89.40%", "-4.51%", "Poor"],
        ["TransKal", "99.09%", "87.57%", "-11.52%", "Severe"],
        ["Conv-AE", "99.39%", "76.04%", "-23.35%", "Failed"],
    ],
    col_widths=[2.5, 2.2, 2.2, 1.8, 2.0]
)

# ============================================================================
# SLIDE 18: Generalization Key Insights
# ============================================================================
add_content_slide(
    "Generalization: Key Insights",
    [
        "CNN-Transformer is the ONLY model that improved on new data!",
        "Transformer attention learns robust, transferable patterns",
        "LSTM and LSTM-FCN show good robustness (<0.5% drop)",
        "TransKal: Kalman parameters overfit to training distribution",
        "Conv-Autoencoder: Threshold highly dataset-specific (needs recalibration)"
    ]
)

# ============================================================================
# SLIDE 19: Final Rankings
# ============================================================================
add_table_slide(
    "Final Model Rankings",
    ["Rank", "Model", "New Data Acc", "Recommendation"],
    [
        ["1", "CNN-Transformer", "99.57%", "Best for deployment"],
        ["2", "LSTM-FCN", "98.93%", "Good alternative"],
        ["3", "LSTM", "98.91%", "Simple, robust"],
        ["4", "XGBoost + HMM", "~91%*", "Interpretable option"],
        ["5", "TransKal", "87.57%", "Avoid (overfit)"],
        ["6", "Conv-AE", "76.04%", "Needs recalibration"],
    ],
    col_widths=[1.5, 3.0, 2.5, 4.5]
)

# ============================================================================
# SLIDE 20: Conclusions
# ============================================================================
add_content_slide(
    "Conclusions",
    [
        "Deep learning significantly outperforms classical ML (99% vs 94%)",
        "CNN-Transformer achieves best generalization (99.57% on new data)",
        "Temporal modeling is critical: sequence-based models excel",
        "HMM filtering helps XGBoost but not deep learning",
        "Autoencoder thresholds require recalibration for new conditions",
        "TransKal's Kalman filter can overfit to training distribution"
    ]
)

# ============================================================================
# SLIDE 21: Practical Recommendations
# ============================================================================
add_content_slide(
    "Practical Recommendations",
    [
        "For multiclass fault classification: CNN-Transformer or LSTM-FCN",
        "For anomaly detection: LSTM Autoencoder (more robust than Conv-AE)",
        "For interpretability: XGBoost with HMM post-processing",
        "Always validate on held-out data from independent conditions",
        "Recalibrate autoencoder thresholds when process conditions change"
    ]
)

# ============================================================================
# SLIDE 22: Limitations & Future Work
# ============================================================================
add_two_column_slide(
    "Limitations & Future Work",
    [
        "Limitations:",
        "  Perfectly balanced (unrealistic)",
        "  No imbalanced evaluation",
        "  17 fault types only",
        "  No novel fault detection",
        "  Sample-level, not time-series metrics"
    ],
    [
        "Future Work:",
        "  Imbalanced test evaluation",
        "  Detection delay metrics",
        "  Novel fault generalization",
        "  Class-wise autoencoder ensemble",
        "  Hierarchical mixture of experts"
    ]
)

# ============================================================================
# SLIDE 23: Acknowledgments
# ============================================================================
add_content_slide(
    "Acknowledgments",
    [
        "Data: Harvard Dataverse (Rieth et al., 2017)",
        "Original TEP: Downs & Vogel (1993)",
        "Computing: Carnegie Mellon University",
        "Contact: jkitchin@andrew.cmu.edu"
    ]
)

# Save the presentation
prs.save('outputs/TEP_ML_Benchmarking_Presentation.pptx')
print("Presentation saved to: outputs/TEP_ML_Benchmarking_Presentation.pptx")
